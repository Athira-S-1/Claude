from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Theme Colors from Reference
NAVY = RGBColor(11, 26, 45)  # #0B1A2D
TEAL = RGBColor(0, 131, 159)  # #00839F
CYAN = RGBColor(29, 141, 255)  # #1D8DFF
GOLD = RGBColor(223, 161, 45)  # #DFA12D
WHITE = RGBColor(255, 255, 255)
LIGHT_GREY = RGBColor(245, 245, 245)
DARK_TEXT = RGBColor(31, 41, 55)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_bg_shape(slide, color):
    """Add background color to slide"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

def add_diagonal_accent(slide):
    """Add diagonal geometric accent like reference"""
    # Navy triangle

    navy_shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(6), Inches(0), Inches(4), Inches(7.5)
    )
    navy_shape.fill.solid()
    navy_shape.fill.fore_color.rgb = NAVY
    navy_shape.line.fill.background()
    navy_shape.rotation = 25
    
def add_title_text(slide, text, top, left, width, height, font_size, color, bold=True):
    """Add formatted title text"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'
    return textbox

def add_body_text(slide, text, top, left, width, height, font_size, color):
    """Add formatted body text"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'
    return textbox


def add_card(slide, left, top, width, height, title, body, color):
    """Add a premium card with shadow"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(230, 230, 230)
    card.line.width = Pt(1)
    card.shadow.inherit = False
    
    # Title
    title_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.2),
        width - Inches(0.4), Inches(0.4)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'
    
    # Body
    body_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.7),
        width - Inches(0.4), height - Inches(0.9)
    )
    tf = body_box.text_frame
    tf.text = body
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Segoe UI'


# ==================== SLIDE 1: COVER PAGE ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg_shape(slide1, WHITE)

# Gold accent line (top left)
gold_line = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.5), Inches(0.6), Inches(0.6), Inches(0.05)
)
gold_line.fill.solid()
gold_line.fill.fore_color.rgb = GOLD
gold_line.line.fill.background()

# Header text
add_body_text(slide1, "CLOUDEX AI SQUAD – KOCHI",
             Inches(0.5), Inches(0.9), Inches(4), Inches(0.3), 12, DARK_TEXT)

# Large ASTRA title
add_title_text(slide1, "ASTRA",
              Inches(1.5), Inches(0.4), Inches(4), Inches(1), 72, NAVY, True)

# Subtitle
add_title_text(slide1, "Autonomous Service Task &\nResolution Assistant",
              Inches(2.5), Inches(0.4), Inches(4.5), Inches(0.8), 24, TEAL, False)

# Description
add_body_text(slide1, "An Embedded Contextual AI Assistant\nfor Enterprise IT Operations",
             Inches(3.3), Inches(0.4), Inches(4.5), Inches(0.5), 16, DARK_TEXT)

# Footer badges
add_body_text(slide1, "📅 Innovation Showcase —\n     Leadership Review",
             Inches(6.3), Inches(0.4), Inches(2.2), Inches(0.4), 11, DARK_TEXT)
add_body_text(slide1, "🎯 Version 2.0",
             Inches(6.3), Inches(2.4), Inches(2), Inches(0.3), 11, DARK_TEXT)

add_body_text(slide1, "📅 August 2026",
             Inches(6.3), Inches(4.4), Inches(2), Inches(0.3), 11, DARK_TEXT)

# Diagonal navy accent (right side)
add_diagonal_accent(slide1)

# Footer bar
footer_bar = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    0, Inches(6.8), prs.slide_width, Inches(0.7)
)
footer_bar.fill.solid()
footer_bar.fill.fore_color.rgb = NAVY
footer_bar.line.fill.background()

# Footer text
footer_text = slide1.shapes.add_textbox(
    Inches(0.7), Inches(6.9), Inches(8), Inches(0.5)
)
tf = footer_text.text_frame
tf.text = "🛡️  Governed. Contextual. Autonomous.\n     Every action. Every record. Every time."
p = tf.paragraphs[0]
p.font.size = Pt(11)
p.font.color.rgb = WHITE
p.font.name = 'Segoe UI'

# ==================== SLIDE 2: EXECUTIVE SUMMARY ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide2, LIGHT_GREY)

# Title
add_title_text(slide2, "Executive Summary",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)

# Subtitle
add_title_text(slide2, "Why Should Leadership Care?",
              Inches(1.1), Inches(0.5), Inches(9), Inches(0.4), 20, TEAL, False)


# Value Proposition
value_box = slide2.shapes.add_textbox(
    Inches(0.5), Inches(1.7), Inches(9), Inches(0.5)
)
tf = value_box.text_frame
tf.text = "ASTRA eliminates 60-70% of manual shift effort while delivering contextual AI assistance that's governed, embedded, and learns from every resolved record."
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = TEAL
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

# KPI Cards
kpi_data = [
    ("60-70%", "Manual Effort\nReduction", NAVY),
    ("< 2 months", "Payback\nPeriod", TEAL),
    ("100+ hrs", "Saved\nper Month", CYAN),
    ("₹1 Crore", "Annual Cost\nAvoidance", GOLD)
]

for i, (metric, label, color) in enumerate(kpi_data):
    left = Inches(0.7 + i * 2.2)
    card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, Inches(2.6), Inches(2), Inches(1.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    # Metric
    metric_box = slide2.shapes.add_textbox(
        left + Inches(0.1), Inches(2.8), Inches(1.8), Inches(0.7)
    )
    tf = metric_box.text_frame
    tf.text = metric
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    
    # Label
    label_box = slide2.shapes.add_textbox(
        left + Inches(0.1), Inches(3.5), Inches(1.8), Inches(0.7)
    )
    tf = label_box.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

# Business Benefits
benefits_title = slide2.shapes.add_textbox(
    Inches(0.5), Inches(4.7), Inches(9), Inches(0.3)
)
tf = benefits_title.text_frame
tf.text = "Key Business Benefits"
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = NAVY

benefits = [
    "✓ Faster incident resolution via instant context assembly",
    "✓ Fewer human errors from repetitive manual data entry",
    "✓ Audit-ready documentation generated as a by-product",
    "✓ Governed foundation that scales safely into automation"
]

for i, benefit in enumerate(benefits):
    benefit_box = slide2.shapes.add_textbox(
        Inches(1), Inches(5.2 + i * 0.35), Inches(8), Inches(0.3)
    )
    tf = benefit_box.text_frame
    tf.text = benefit
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT

# The Ask
ask_box = slide2.shapes.add_textbox(
    Inches(0.5), Inches(6.7), Inches(9), Inches(0.6)
)

tf = ask_box.text_frame
tf.text = "The Ask: Approval for Phase 1 Foundation (Weeks 1-8) — sandboxed, low-risk validation with measurable outcomes."
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEAL
p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 3: WHAT IS ASTRA? ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide3, WHITE)

# Title
add_title_text(slide3, "What Is ASTRA?",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)

# Subtitle
quote_box = slide3.shapes.add_textbox(
    Inches(1), Inches(1.2), Inches(8), Inches(0.5)
)
tf = quote_box.text_frame
tf.text = '"What would you like me to do?" — asked only after ASTRA already knows what\'s going on.'
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = TEAL
p.alignment = PP_ALIGN.CENTER

# Workflow Cards
workflow = [
    ("1", "UNDERSTAND", "Identify the active\nrecord & lifecycle stage", NAVY),
    ("2", "ANALYSE", "Assemble unified\ncontext", TEAL),
    ("3", "RECOMMEND", "Surface the most\nrelevant actions", NAVY),
    ("4", "ASSIST", "Support investigation\n& follow-up", CYAN),
    ("5", "DOCUMENT", "Generate notes as a\nby-product of work", NAVY)
]


for i, (num, title, desc, color) in enumerate(workflow):
    left = Inches(0.5 + i * 1.85)
    
    # Card
    card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, Inches(2.3), Inches(1.7), Inches(2.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    
    # Number circle
    circle = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Inches(0.6), Inches(2.5), Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    
    num_box = slide3.shapes.add_textbox(
        left + Inches(0.6), Inches(2.5), Inches(0.5), Inches(0.5)
    )
    tf = num_box.text_frame
    tf.text = num
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Title
    title_box = slide3.shapes.add_textbox(
        left + Inches(0.1), Inches(3.2), Inches(1.5), Inches(0.4)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    
    # Description
    desc_box = slide3.shapes.add_textbox(
        left + Inches(0.1), Inches(3.7), Inches(1.5), Inches(1.2)
    )
    tf = desc_box.text_frame
    tf.text = desc
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Arrow (except last card)
    if i < len(workflow) - 1:
        arrow = slide3.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            left + Inches(1.75), Inches(3.5), Inches(0.4), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD if i % 2 == 1 else CYAN
        arrow.line.fill.background()

# Feedback loop note
feedback_box = slide3.shapes.add_textbox(
    Inches(1), Inches(5.6), Inches(8), Inches(0.6)
)
tf = feedback_box.text_frame
tf.text = "↩ Feedback loop — every closed record enriches the knowledge used at Stage 2 for the next event"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = TEAL
p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 4: CURRENT STATE ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide4, LIGHT_GREY)


# Title
add_title_text(slide4, "Current State",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)
add_title_text(slide4, "What Happens Today?",
              Inches(1.1), Inches(0.5), Inches(9), Inches(0.4), 20, TEAL, False)

# Manual workflow
manual_steps = [
    ("📋", "Engineer\nOpens Ticket"),
    ("🔍", "Searches\n5 Systems"),
    ("📝", "Copies Data\nManually"),
    ("🤔", "Reconstructs\nContext"),
    ("⏱️", "30+ min\nBefore Action")
]

for i, (icon, label) in enumerate(manual_steps):
    left = Inches(0.8 + i * 1.7)
    
    # Icon box
    icon_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, Inches(2.5), Inches(1.3), Inches(1.3)
    )
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = WHITE
    icon_box.line.color.rgb = NAVY
    icon_box.line.width = Pt(2)
    
    # Icon
    icon_text = slide4.shapes.add_textbox(
        left + Inches(0.2), Inches(2.7), Inches(0.9), Inches(0.5)
    )
    tf = icon_text.text_frame
    tf.text = icon
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    label_text = slide4.shapes.add_textbox(
        left, Inches(3.3), Inches(1.3), Inches(0.5)
    )

    tf = label_text.text_frame
    tf.text = label
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Arrow
    if i < len(manual_steps) - 1:
        arrow = slide4.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            left + Inches(1.35), Inches(3), Inches(0.35), Inches(0.15)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

# Impact statement
impact = slide4.shapes.add_textbox(
    Inches(1), Inches(5), Inches(8), Inches(1.5)
)
tf = impact.text_frame
tf.text = "Result:\n• 3 engineers per shift on manual context reconstruction\n• 100+ hours/month lost to repetitive data gathering\n• Diagnosis restarts from zero for already-seen patterns"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = DARK_TEXT
p.line_spacing = 1.5

# ==================== SLIDE 5: PAIN POINTS ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide5, WHITE)

# Title
add_title_text(slide5, "Pain Points",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)

add_title_text(slide5, "What Problem Are We Solving?",
              Inches(1.1), Inches(0.5), Inches(9), Inches(0.4), 20, TEAL, False)

# Pain point cards
pain_points = [
    ("HRS", "High Man-Hours", "3 engineers/shift spent on\nmanual context reconstruction",
     RGBColor(185, 28, 28)),  # Red
    ("!", "High Human Error", "Manual, repetitive data entry\nacross disconnected systems",
     RGBColor(217, 119, 6)),  # Orange
    ("TIME", "High Resolution Time", "Diagnosis restarts from zero\nfor already-seen patterns",
     NAVY)
]

for i, (badge, title, desc, color) in enumerate(pain_points):
    left = Inches(0.8 + i * 3)
    
    # Card
    card = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, Inches(2), Inches(2.6), Inches(3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    
    # Badge circle
    badge_circle = slide5.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Inches(0.8), Inches(2.4), Inches(1), Inches(1)
    )
    badge_circle.fill.solid()
    badge_circle.fill.fore_color.rgb = RGBColor(255, 255, 255, alpha=0.3)
    badge_circle.line.fill.background()

    
    # Badge text
    badge_text = slide5.shapes.add_textbox(
        left + Inches(0.8), Inches(2.5), Inches(1), Inches(0.8)
    )
    tf = badge_text.text_frame
    tf.text = badge
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Title
    title_box = slide5.shapes.add_textbox(
        left + Inches(0.2), Inches(3.6), Inches(2.2), Inches(0.4)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide5.shapes.add_textbox(
        left + Inches(0.2), Inches(4.1), Inches(2.2), Inches(0.7)
    )
    tf = desc_box.text_frame
    tf.text = desc
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Impact summary
summary_box = slide5.shapes.add_textbox(
    Inches(1), Inches(5.7), Inches(8), Inches(0.8)
)

tf = summary_box.text_frame
tf.text = "Combined Impact: 100+ hrs/month for incident context, 80-120 hrs/month for change & problem review, 20-60 hrs/month for governance & operational drift across representative accounts."
tf.word_wrap = True
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = TEAL
p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 6: BUSINESS CONSEQUENCES ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide6, LIGHT_GREY)

# Title
add_title_text(slide6, "Business Consequences",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)
add_title_text(slide6, "Why Act Now?",
              Inches(1.1), Inches(0.5), Inches(9), Inches(0.4), 20, TEAL, False)

# KPI Impact Dashboard
kpis = [
    ("200-300 hrs", "Manual Effort\nper Month", NAVY),
    ("₹99.6L", "Annual Cost\nAvoidance", GOLD),
    ("3:1", "Engineer to\nASTRA Ratio", TEAL)
]

for i, (value, label, color) in enumerate(kpis):
    left = Inches(1.3 + i * 2.5)
    
    # Card
    card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, Inches(2.2), Inches(2.2), Inches(1.6)
    )

    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(3)
    
    # Value
    value_box = slide6.shapes.add_textbox(
        left + Inches(0.1), Inches(2.4), Inches(2), Inches(0.7)
    )
    tf = value_box.text_frame
    tf.text = value
    p = tf.paragraphs[0]
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    label_box = slide6.shapes.add_textbox(
        left + Inches(0.1), Inches(3.1), Inches(2), Inches(0.5)
    )
    tf = label_box.text_frame
    tf.text = label
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

# Business risks
risks_title = slide6.shapes.add_textbox(
    Inches(0.5), Inches(4.2), Inches(9), Inches(0.4)
)
tf = risks_title.text_frame
tf.text = "Business Risks of Inaction"
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = NAVY

risks = [
    "⚠️ Operational inefficiency compounds as service portfolio grows",

    "⚠️ Knowledge loss accelerates with engineer turnover",
    "⚠️ Repetitive manual effort increases human error & burnout",
    "⚠️ Competitors deploy AI-assisted operations, increasing competitive pressure"
]

for i, risk in enumerate(risks):
    risk_box = slide6.shapes.add_textbox(
        Inches(1), Inches(4.7 + i * 0.4), Inches(8), Inches(0.35)
    )
    tf = risk_box.text_frame
    tf.text = risk
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT

# ==================== SLIDE 7: PRODUCT DIFFERENTIATORS ====================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide7, WHITE)

# Title
add_title_text(slide7, "Product Differentiators",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)
add_title_text(slide7, "Why ASTRA?",
              Inches(1.1), Inches(0.5), Inches(9), Inches(0.4), 20, TEAL, False)

# Comparison matrix header
headers = ["Capability", "Generic AI\nChatbots", "Static ITSM\nDashboards", "ASTRA"]
header_colors = [DARK_TEXT, DARK_TEXT, DARK_TEXT, TEAL]

for i, (header, color) in enumerate(zip(headers, header_colors)):
    left = Inches(0.7 + i * 2.2)
    width = Inches(2) if i == 0 else Inches(1.8)

    
    header_box = slide7.shapes.add_textbox(
        left, Inches(2), width, Inches(0.5)
    )
    tf = header_box.text_frame
    tf.text = header
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

# Comparison rows
comparisons = [
    ("Understands active\nrecord automatically", "➖", "➖", "✓"),
    ("Embedded in ITSM\nworkspace", "➖", "✓", "✓"),
    ("Explains its\nreasoning", "➖", "➖", "✓"),
    ("Governed tiered\nautonomy", "➖", "➖", "✓"),
    ("Learns from every\nresolved record", "➖", "➖", "✓")
]

for row_idx, (capability, generic, static, astra) in enumerate(comparisons):
    top = Inches(2.7 + row_idx * 0.7)
    
    # Capability
    cap_box = slide7.shapes.add_textbox(
        Inches(0.7), top, Inches(2), Inches(0.6)
    )
    tf = cap_box.text_frame
    tf.text = capability
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT
    
    # Columns
    for col_idx, value in enumerate([generic, static, astra]):
        left = Inches(2.9 + col_idx * 2.2)

        
        val_box = slide7.shapes.add_textbox(
            left, top + Inches(0.05), Inches(1.8), Inches(0.5)
        )
        tf = val_box.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEAL if value == "✓" else RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 8: TECHNOLOGIES USED ====================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide8, LIGHT_GREY)

# Title
add_title_text(slide8, "Technologies Used",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)
add_title_text(slide8, "How Is It Built?",
              Inches(1.1), Inches(0.5), Inches(9), Inches(0.4), 20, TEAL, False)

# Layered architecture
layers = [
    ("User Experience Layer", "Embedded Contextual AI Assistant · ServiceNow", NAVY),
    ("Context Intelligence Layer", "User · Record · Business · Technical · Historical · Knowledge context", TEAL),
    ("AI Orchestration & Reasoning", "Intent detection · Skill/Agent selection · Governance checks", CYAN),
    ("AI Agents & Enterprise Skills", "LLM-grounded Agents · Enterprise-governed Skills", RGBColor(0, 131, 159)),
    ("Enterprise Knowledge & RAG", "Retrieval-Augmented Generation over organizational knowledge", RGBColor(217, 119, 6)),
    ("Integration, Security & Governance", "ITSM · Cloud · Monitoring · RBAC · Identity · Audit trail", RGBColor(120, 53, 15))
]

for i, (layer_name, layer_desc, color) in enumerate(layers):
    top = Inches(2 + i * 0.7)
    
    # Layer bar
    layer_bar = slide8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), top, Inches(8), Inches(0.6)
    )
    layer_bar.fill.solid()
    layer_bar.fill.fore_color.rgb = color
    layer_bar.line.fill.background()
    
    # Layer name
    name_box = slide8.shapes.add_textbox(
        Inches(1.2), top + Inches(0.05), Inches(7.6), Inches(0.25)
    )
    tf = name_box.text_frame
    tf.text = layer_name
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Layer description
    desc_box = slide8.shapes.add_textbox(
        Inches(1.2), top + Inches(0.3), Inches(7.6), Inches(0.25)
    )
    tf = desc_box.text_frame
    tf.text = layer_desc
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE

# Technology stack
tech_title = slide8.shapes.add_textbox(
    Inches(1), Inches(6.2), Inches(8), Inches(0.3)
)

tf = tech_title.text_frame
tf.text = "Core Technologies: Azure OpenAI / Claude AI · Semantic Kernel · Vector DBs · ServiceNow/Jira/BMC · Existing Runbooks"
p = tf.paragraphs[0]
p.font.size = Pt(12)
p.font.color.rgb = DARK_TEXT
p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 9: COMPLETED FUNCTIONALITIES ====================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide9, WHITE)

# Title
add_title_text(slide9, "Completed Functionalities",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)
add_title_text(slide9, "What Is Already Working?",
              Inches(1.1), Inches(0.5), Inches(9), Inches(0.4), 20, TEAL, False)

# Capability cards (2x2 grid)
capabilities = [
    ("🎯 Proactive Certificate\n& Cost Hygiene", 
     "Scheduled scans detect expiring certificates & orphaned resources; raises proactive ticket with action plan",
     "✅ Working", TEAL),
    ("📡 Service-Health\nWebhook Auto-Triage",
     "Cloud service-health events auto-create scoped incident with draft action plan for engineer review",
     "✅ Working", CYAN),
    ("📋 Incomplete Ticket\nData Enforcement",
     "Active field prompts request specific data when required fields missing; prevents incomplete ticket closure",
     "✅ Working", TEAL),
    ("🎯 Stakeholder &\nAssignment Routing",
     "Matches issue pattern against ownership/skill data; suggests correct assignment group & escalation path",
     "✅ Working", CYAN)
]

positions = [
    (Inches(0.7), Inches(2.2)),
    (Inches(5.1), Inches(2.2)),
    (Inches(0.7), Inches(4.6)),
    (Inches(5.1), Inches(4.6))
]

for (left, top), (title, desc, status, color) in zip(positions, capabilities):
    # Card
    card = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, Inches(4.2), Inches(2.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    # Title
    title_box = slide9.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.2), Inches(3.8), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    
    # Description
    desc_box = slide9.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.8), Inches(3.8), Inches(1)
    )

    tf = desc_box.text_frame
    tf.text = desc
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.line_spacing = 1.3
    
    # Status badge
    status_badge = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(3), top + Inches(1.8), Inches(1), Inches(0.3)
    )
    status_badge.fill.solid()
    status_badge.fill.fore_color.rgb = color
    status_badge.line.fill.background()
    
    status_text = slide9.shapes.add_textbox(
        left + Inches(3), top + Inches(1.8), Inches(1), Inches(0.3)
    )
    tf = status_text.text_frame
    tf.text = status
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.vertical_anchor = MSO_ANCHOR.MIDDLE

# ==================== SLIDE 10: PROTOTYPE ====================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide10, LIGHT_GREY)

# Title
add_title_text(slide10, "Prototype",
              Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), 40, NAVY)
add_title_text(slide10, "Can We See It?",
              Inches(1.1), Inches(0.5), Inches(9), Inches(0.4), 20, TEAL, False)


# Prototype screenshot placeholder
screenshot_box = slide10.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1), Inches(2), Inches(8), Inches(3.8)
)
screenshot_box.fill.solid()
screenshot_box.fill.fore_color.rgb = WHITE
screenshot_box.line.color.rgb = NAVY
screenshot_box.line.width = Pt(2)

# Screenshot placeholder text
screenshot_text = slide10.shapes.add_textbox(
    Inches(2), Inches(3), Inches(6), Inches(2)
)
tf = screenshot_text.text_frame
tf.text = "🖥️\n\nASTRA-Embedded Work Queue\n\nServiceNow (sandbox) · Incident · Problem · Change\n\nContextual AI Assistant with:\n✓ Automatic context assembly\n✓ Recommended actions\n✓ Draft handover notes\n✓ Investigation assistance"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = DARK_TEXT
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.5

# Live demo callout
demo_box = slide10.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.5), Inches(6), Inches(3), Inches(0.6)
)
demo_box.fill.solid()
demo_box.fill.fore_color.rgb = GOLD
demo_box.line.fill.background()

demo_text = slide10.shapes.add_textbox(
    Inches(3.5), Inches(6), Inches(3), Inches(0.6)
)

tf = demo_text.text_frame
tf.text = "🎬  Live Demo Available"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p.vertical_anchor = MSO_ANCHOR.MIDDLE

# Workflow annotation
workflow_note = slide10.shapes.add_textbox(
    Inches(1), Inches(6.8), Inches(8), Inches(0.5)
)
tf = workflow_note.text_frame
tf.text = "Annotated Workflow: Engineer opens record → ASTRA assembles context → suggests actions (Tier 2) → engineer approves → work notes + KB article generated"
tf.word_wrap = True
p = tf.paragraphs[0]
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = TEAL
p.alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('/projects/sandbox/ASTRA-Leadership-Review.pptx')
print("✅ PowerPoint created successfully: ASTRA-Leadership-Review.pptx")
print("\n📋 Slide Summary:")
print("1. Cover Page - Professional theme with ASTRA branding")
print("2. Executive Summary - KPI cards & business benefits")
print("3. What Is ASTRA? - 5-step workflow with visual cards")
print("4. Current State - Manual workflow visualization")
print("5. Pain Points - 3 premium problem cards")
print("6. Business Consequences - Impact dashboard & risks")
print("7. Product Differentiators - Comparison matrix")
print("8. Technologies Used - Layered architecture diagram")
print("9. Completed Functionalities - 4 capability cards")
print("10. Prototype - Screenshot placeholder & live demo CTA")
